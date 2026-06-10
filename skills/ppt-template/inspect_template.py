#!/usr/bin/env python3
"""
inspect_template.py — 새 PPT 템플릿을 build_deck.py 에 맞추기 위한 분석 도구.

새 template.pptx 를 받으면 이 스크립트로 (1)레이아웃 이름 (2)각 레이아웃·샘플
슬라이드의 도형 좌표/폰트 를 덤프한다. 그 출력을 보고 build_deck.py 상단의
LAYOUTS/폰트/색상과 build_* 함수의 좌표를 새 값으로 맞춘다. (PORTING.md 참조)

사용법:
  python3 inspect_template.py template.pptx              # 레이아웃 + 샘플슬라이드 요약
  python3 inspect_template.py template.pptx --layouts    # 레이아웃 도형까지 상세
  python3 inspect_template.py template.pptx --slide 8    # 특정 샘플 슬라이드 상세
"""
import argparse
from pptx import Presentation
from pptx.util import Emu


def inch(v):
    return round(Emu(v).inches, 2) if v is not None else None


def shape_line(sh):
    t = ""
    fs = ""
    if sh.has_text_frame:
        t = sh.text_frame.text[:60].replace("\n", " / ")
        for p in sh.text_frame.paragraphs:
            if p.runs:
                r = p.runs[0]
                col = ""
                try:
                    col = str(r.font.color.rgb)
                except Exception:
                    pass
                fs = (f" font={r.font.name},"
                      f"{r.font.size.pt if r.font.size else '?'}pt,"
                      f"bold={r.font.bold},color={col}")
                break
    ph = ""
    try:
        if sh.is_placeholder:
            ph = f" PH(idx={sh.placeholder_format.idx},type={sh.placeholder_format.type})"
    except Exception:
        pass
    return (f"  [{sh.shape_type}] {sh.name!r}{ph} "
            f"pos=({inch(sh.left)},{inch(sh.top)}) "
            f"size=({inch(sh.width)}x{inch(sh.height)}){fs} text={t!r}")


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("template")
    ap.add_argument("--layouts", action="store_true",
                    help="각 레이아웃의 도형까지 상세 출력")
    ap.add_argument("--slide", type=int, default=None,
                    help="특정 샘플 슬라이드(0-base) 상세 출력")
    args = ap.parse_args()

    p = Presentation(args.template)
    print(f"슬라이드 크기: {inch(p.slide_width)} x {inch(p.slide_height)} in")
    print(f"샘플 슬라이드 수: {len(p.slides._sldIdLst)}")

    print("\n=== 레이아웃 목록 (이것을 build_deck.py 의 LAYOUTS 에 매핑) ===")
    for mi, master in enumerate(p.slide_masters):
        for li, lay in enumerate(master.slide_layouts):
            nph = len(list(lay.placeholders))
            print(f"  Master[{mi}] Layout[{li}]: {lay.name!r}  (placeholders={nph})")
            if args.layouts:
                for sh in lay.shapes:
                    print("    " + shape_line(sh))

    if args.slide is not None:
        s = p.slides[args.slide]
        print(f"\n=== 샘플 슬라이드[{args.slide}] layout={s.slide_layout.name!r} ===")
        for sh in s.shapes:
            print(shape_line(sh))
    else:
        print("\n=== 샘플 슬라이드 레이아웃 분포 ===")
        from collections import Counter
        c = Counter(s.slide_layout.name for s in p.slides)
        for name, n in c.items():
            print(f"  {name!r}: {n}장")
        print("\n팁: 본문 슬라이드의 좌표를 보려면 --slide <번호> 로 상세 확인하세요.")


if __name__ == "__main__":
    main()
