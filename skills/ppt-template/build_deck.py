#!/usr/bin/env python3
"""
build_deck.py — 사내 보고서 템플릿(보고서양식.pptx / POSCO IH 가로형)을 베이스로
PPTX 보고서를 생성한다.

핵심 원칙
  - 절대 빈 deck에서 시작하지 않는다. template.pptx 를 열어 디자인(배경 그래픽이
    들어있는 슬라이드 레이아웃)을 그대로 상속한 뒤, 정해진 좌표에 텍스트박스만 올린다.
  - 좌표/폰트/크기는 실제 템플릿 슬라이드에서 추출한 값을 따른다.
  - 본문은 항상 '본문1'(빈 레이아웃)을 베이스로 쓰고, 단/표/이미지 배치는 코드로
    구성한다. (본문2/3/4 레이아웃에는 빨간 가이드 박스가 배경에 박혀 있어 사용 안 함)

사용법
  python3 build_deck.py spec.json -o output.pptx
  python3 build_deck.py spec.json            # -> output.pptx

spec.json 스키마는 같은 폴더의 SKILL.md 와 spec.schema.json 을 참조.
"""

import argparse
import json
import os
import re
import sys

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

HERE = os.path.dirname(os.path.abspath(__file__))
TEMPLATE = os.path.join(HERE, "template.pptx")

# spec 파일이 있는 폴더 (이미지 상대경로 해석 기준). build()에서 설정.
SPEC_DIR = os.getcwd()

# 슬라이드 크기(이 템플릿): 10.83 x 7.5 in
SLIDE_W = 10.83


def _resolve(path, sec=None):
    """이미지 경로 해석: 절대경로면 그대로, 상대경로면 spec 폴더 기준."""
    if os.path.isabs(path) or os.path.exists(path):
        return path
    cand = os.path.join(SPEC_DIR, path)
    if not os.path.exists(cand):
        raise FileNotFoundError(f"이미지를 찾을 수 없습니다: {path} (기준: {SPEC_DIR})")
    return cand

# ====================================================================
# 새 템플릿으로 교체할 때는 아래 블록만 수정하면 된다 (PORTING.md 참조).
# inspect_template.py 로 새 template.pptx 의 레이아웃 이름/좌표/폰트를 먼저 확인할 것.
# --------------------------------------------------------------------
# (1) 레이아웃 이름 — 새 템플릿의 슬라이드 레이아웃 이름과 1:1로 맞춘다.
#     이 템플릿은 목차/간지/엔딩이 모두 '간지' 레이아웃을 재사용한다.
LAYOUTS = {
    "cover":    "표지",
    "foreword": "들어가며",
    "toc":      "간지",
    "divider":  "간지",
    "body":     "본문1",
    "ending":   "간지",
}

# (2) 폰트 패밀리 — 새 템플릿의 사내 폰트로 바꾼다.
F_TITLE = "HY견고딕"    # 상단 장 라벨 / 표지 / 목차 / 간지 제목
F_HEAD = "맑은 고딕"    # 본문 대제목 22pt / ■ 소제목 18pt (bold)
F_BODY = "맑은 고딕"    # 본문 11pt / 표
F_SRC = "맑은 고딕"     # 출처 9pt
F_DIV = "HY견고딕"      # 간지 대제목 32pt / 엔딩 48pt
F_FORE = "바탕"         # 들어가며(서문) 본문 14pt

# (3) 색상 — 새 템플릿의 브랜드 컬러로 바꾼다.
NAVY = RGBColor(0x1F, 0x49, 0x7D)   # 표지 제목 / 엔딩
BLUE = RGBColor(0x25, 0x3C, 0x9F)   # 본문 대제목(밝은 네이비)
HEADER = RGBColor(0x44, 0x54, 0x6A) # 상단 장 라벨 바
DARK = RGBColor(0x17, 0x37, 0x5E)   # 목차/간지 제목 (짙은 네이비)
INK = RGBColor(0x16, 0x24, 0x3E)    # 표지 날짜/기관 (매우 짙은 네이비)
GRAY = RGBColor(0x59, 0x59, 0x59)   # 본문 텍스트
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
# (4) 각 슬라이드 요소의 좌표(Inches)는 아래 build_* 함수 안에 있다.
#     좌표가 어긋나면 render_preview.sh 로 렌더해 눈으로 보며 조정한다.
# ====================================================================

# ---------------------------------------------------------------- 강조(하이라이트) 마크업
# 본문/표/head 텍스트에 인라인 마크업으로 핵심을 강조한다.
#   **키워드**  -> 네이비 볼드 (핵심 용어 강조)
#   ==수치==    -> 레드 볼드 + 노란 형광 (핵심 수치/지표 하이라이트)
# 마크업이 없으면 단일 run과 동일하게 동작(기존 spec 호환).
HL_TERM = BLUE                       # 키워드 강조색(네이비)
HL_NUM = RGBColor(0xC0, 0x39, 0x2B)  # 수치 강조색(레드)
HL_BG = "FFF29A"                     # 수치 형광(노랑)
_MARK_RE = re.compile(r"(\*\*.+?\*\*|==.+?==)")


def _set_run_highlight(run, hexstr):
    """run 에 PowerPoint 텍스트 형광(a:highlight)을 입힌다(지원 환경에서 표시)."""
    rPr = run._r.get_or_add_rPr()
    for old in rPr.findall(qn("a:highlight")):
        rPr.remove(old)
    hl = rPr.makeelement(qn("a:highlight"), {})
    hl.append(rPr.makeelement(qn("a:srgbClr"), {"val": hexstr}))
    ref = rPr.find(qn("a:solidFill"))
    if ref is not None:
        ref.addnext(hl)
    else:
        rPr.insert(0, hl)


def _emit_runs(p, text, font, size, color, bold):
    """text 를 **키워드** / ==수치== 마크업 단위로 분해해 paragraph에 run으로 추가."""
    for tok in _MARK_RE.split(text):
        if tok == "":
            continue
        kind = None
        if len(tok) >= 4 and tok.startswith("**") and tok.endswith("**"):
            tok, kind = tok[2:-2], "term"
        elif len(tok) >= 4 and tok.startswith("==") and tok.endswith("=="):
            tok, kind = tok[2:-2], "num"
        run = p.add_run()
        run.text = tok
        f = run.font
        f.name = font
        f.size = Pt(size)
        if kind == "term":
            f.bold = True
            f.color.rgb = HL_TERM
        elif kind == "num":
            f.bold = True
            f.color.rgb = HL_NUM
            _set_run_highlight(run, HL_BG)
        else:
            f.bold = bold
            f.color.rgb = color
    return p


# ---------------------------------------------------------------- helpers
def layout_by_name(prs, name):
    for master in prs.slide_masters:
        for lay in master.slide_layouts:
            if lay.name == name:
                return lay
    raise KeyError(f"레이아웃 '{name}' 을(를) 템플릿에서 찾을 수 없습니다.")


def wipe_slides(prs):
    """템플릿에 포함된 샘플 슬라이드를 모두 제거(레이아웃/마스터는 유지).

    슬라이드 ID 목록뿐 아니라 부모 파트의 관계(rId)까지 끊어 슬라이드 파트가
    저장 시 패키지에서 빠지도록 한다. (안 하면 slideN.xml 이름이 충돌함)
    """
    xml_slides = prs.slides._sldIdLst
    for sldId in list(xml_slides):
        prs.part.drop_rel(sldId.rId)
        xml_slides.remove(sldId)


def add_slide(prs, layout_name):
    return prs.slides.add_slide(layout_by_name(prs, layout_name))


def textbox(slide, left, top, width, height, text,
            font=F_BODY, size=11, bold=False, color=GRAY,
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.0,
            wrap=True, rich=False):
    """단일 문자열 텍스트박스. wrap=False면 짧은 라벨이 줄바꿈되지 않음.
    rich=True면 **키워드**/==수치== 마크업을 파싱해 강조 run을 만든다."""
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                  Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    if rich:
        _emit_runs(p, text, font, size, color, bold)
    else:
        run = p.add_run()
        run.text = text
        f = run.font
        f.name = font
        f.size = Pt(size)
        f.bold = bold
        f.color.rgb = color
    return tb


def paragraphs(slide, left, top, width, height, lines,
               font=F_BODY, size=11, bold=False, color=GRAY,
               align=PP_ALIGN.LEFT, line_spacing=1.2, space_after=6):
    """여러 단락 텍스트박스. lines 는 문자열 리스트(글머리표 없음)."""
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                  Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        _emit_runs(p, line, font, size, color, bold)
    return tb


def bullets(slide, left, top, width, height, items,
            font=F_BODY, size=11, color=GRAY, bullet_char="• ",
            bold=False, line_spacing=1.15, space_after=6):
    """여러 줄 글머리표 텍스트박스. items 는 문자열 리스트."""
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                  Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        # 글머리표 기호(강조 없음) + 본문(마크업 파싱)
        r0 = p.add_run()
        r0.text = bullet_char
        f0 = r0.font
        f0.name = font
        f0.size = Pt(size)
        f0.bold = bold
        f0.color.rgb = color
        _emit_runs(p, item, font, size, color, bold)
    return tb


def hline(slide, left, top, width, color=NAVY, weight=1.25):
    from pptx.enum.shapes import MSO_CONNECTOR
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                    Inches(left), Inches(top),
                                    Inches(left + width), Inches(top))
    ln.line.color.rgb = color
    ln.line.width = Pt(weight)
    return ln


def vline(slide, left, top, height, color=HEADER, weight=1.0):
    from pptx.enum.shapes import MSO_CONNECTOR
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                    Inches(left), Inches(top),
                                    Inches(left), Inches(top + height))
    ln.line.color.rgb = color
    ln.line.width = Pt(weight)
    return ln


def _set_cell(cell, text, size=10, bold=False, color=GRAY,
              fill=None, align=PP_ALIGN.LEFT, rich=False):
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = Inches(0.06)
    cell.margin_right = Inches(0.06)
    cell.margin_top = Inches(0.02)
    cell.margin_bottom = Inches(0.02)
    if fill is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
    else:
        cell.fill.background()
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    if rich:
        _emit_runs(p, str(text), F_BODY, size, color, bold)
    else:
        run = p.add_run()
        run.text = str(text)
        f = run.font
        f.name = F_BODY
        f.size = Pt(size)
        f.bold = bold
        f.color.rgb = color


def add_table(slide, left, top, width, height, headers, rows,
              size=10):
    """네이비 헤더 + 흰 본문의 깔끔한 표. headers 없으면 헤더행 생략."""
    has_header = bool(headers)
    ncols = len(headers) if has_header else (len(rows[0]) if rows else 1)
    nrows = len(rows) + (1 if has_header else 0)
    nrows = max(nrows, 1)
    gframe = slide.shapes.add_table(nrows, ncols, Inches(left), Inches(top),
                                    Inches(width), Inches(height))
    tbl = gframe.table
    # 기본 줄무늬 스타일 끄기 (셀 채움을 직접 제어)
    tbl.first_row = has_header
    tbl.horz_banding = False
    r0 = 0
    if has_header:
        for c, h in enumerate(headers):
            _set_cell(tbl.cell(0, c), h, size=size, bold=True,
                      color=WHITE, fill=NAVY, align=PP_ALIGN.CENTER)
        r0 = 1
    for ri, row in enumerate(rows):
        for c in range(ncols):
            val = row[c] if c < len(row) else ""
            _set_cell(tbl.cell(r0 + ri, c), val, size=size, color=GRAY,
                      fill=WHITE, rich=True)
    return gframe


def add_image_fit(slide, path, left, top, max_w, max_h):
    """비율 유지하며 (max_w x max_h) 박스에 맞춰 이미지 배치(좌상단 기준)."""
    pic = slide.shapes.add_picture(path, Inches(left), Inches(top),
                                   width=Inches(max_w))
    if pic.height > Inches(max_h):
        ratio = Inches(max_h) / pic.height
        pic.width = int(pic.width * ratio)
        pic.height = int(pic.height * ratio)
    return pic


# ---------------------------------------------------------------- slide builders
def build_cover(prs, spec):
    """표지. 모든 요소는 슬라이드 전체 폭에 가운데 정렬."""
    s = add_slide(prs, LAYOUTS["cover"])
    # 보고서 종류 배지 (예: 완료보고서)
    rtype = spec.get("report_type", "완료보고서")
    if rtype:
        textbox(s, 0.0, 1.33, SLIDE_W, 0.45, rtype,
                font=F_TITLE, size=20, color=INK, align=PP_ALIGN.CENTER)
    # 제목
    title = spec.get("title", "제목 없음")
    textbox(s, 0.0, 1.97, SLIDE_W, 1.2, title,
            font=F_TITLE, size=40, bold=False, color=NAVY,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP, line_spacing=1.05)
    # 부제 (선택)
    sub = spec.get("subtitle")
    if sub:
        textbox(s, 0.0, 5.45, SLIDE_W, 0.5, sub,
                font=F_HEAD, size=15, color=GRAY, align=PP_ALIGN.CENTER)
    # 하단: 날짜 / 기관
    date = spec.get("date")
    if date:
        textbox(s, 0.0, 6.17, SLIDE_W, 0.38, date,
                font=F_HEAD, size=18, bold=True, color=INK,
                align=PP_ALIGN.CENTER)
    org = spec.get("org", "POSCO IH")
    if org:
        textbox(s, 0.0, 6.65, SLIDE_W, 0.38, org,
                font=F_HEAD, size=20, bold=True, color=INK,
                align=PP_ALIGN.CENTER)


def build_foreword(prs, spec):
    """들어가며(서문). 발주처 수신·제목 헤더 + 서문 본문."""
    fw = spec.get("foreword")
    if not fw:
        return
    s = add_slide(prs, LAYOUTS["foreword"])
    # 상단 헤더(수신/제목) — 16pt bold
    head_lines = []
    if fw.get("recipient"):
        head_lines.append(f"수신 : {fw['recipient']}")
    if fw.get("title"):
        head_lines.append(f"제목 : {fw['title']}")
    if head_lines:
        paragraphs(s, 0.45, 0.45, 9.9, 1.3, head_lines,
                   font=F_HEAD, size=16, bold=True, color=INK,
                   line_spacing=1.3, space_after=8)
    # 서문 본문 — 14pt 바탕. body 는 문자열(줄바꿈 분리) 또는 리스트.
    body = fw.get("body", "")
    if isinstance(body, str):
        lines = [ln for ln in body.split("\n") if ln.strip()]
    else:
        lines = list(body)
    if lines:
        paragraphs(s, 0.65, 3.0, 9.53, 2.6, lines,
                   font=F_FORE, size=14, color=GRAY,
                   line_spacing=1.4, space_after=10)


def build_toc(prs, spec):
    items = spec.get("toc") or [sec.get("divider") or sec.get("section_title")
                                for sec in spec.get("sections", [])
                                if sec.get("divider") or sec.get("section_title")]
    if not items:
        return
    s = add_slide(prs, LAYOUTS["toc"])
    textbox(s, 1.06, 0.9, 6, 0.7, "목차", font=F_TITLE, size=28,
            bold=False, color=DARK)
    top = 1.95
    # 항목이 많으면 간격을 자동으로 좁혀 한 슬라이드에 맞춘다(기본 0.78").
    step = min(0.78, (7.15 - top) / max(len(items), 1))
    for i, it in enumerate(items, 1):
        textbox(s, 1.3, top, 1.0, 0.5, f"{i:02d}",
                font=F_TITLE, size=20, bold=True, color=BLUE, wrap=False)
        textbox(s, 2.3, top + 0.04, 8.0, 0.5, it,
                font=F_HEAD, size=17, color=DARK)
        top += step


def build_divider(prs, text):
    """간지(장 구분). 짙은 네이비 대제목, 페이지 가운데."""
    s = add_slide(prs, LAYOUTS["divider"])
    textbox(s, 0.0, 3.12, SLIDE_W, 0.9, text,
            font=F_DIV, size=32, bold=False, color=DARK,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP, line_spacing=1.05)


def _body_header(s, sec):
    """본문 공통 상단부(장 라벨/대제목/소제목)를 그리고 콘텐츠 시작 top을 반환."""
    # 상단 장 라벨 바
    if sec.get("section_title"):
        textbox(s, 0.14, 0.14, 10.55, 0.43, sec["section_title"],
                font=F_TITLE, size=20, color=HEADER)
    # 대제목 = head (핵심 한 줄), 가운데 정렬 + 밑줄
    head = sec.get("head", "")
    if head:
        textbox(s, 0.3, 0.82, 10.2, 0.55, head,
                font=F_HEAD, size=22, bold=True, color=BLUE,
                align=PP_ALIGN.CENTER, line_spacing=1.0, rich=True)
        hline(s, 0.67, 1.42, 9.57, color=NAVY, weight=1.25)
    # ■ 소제목
    content_top = 1.7
    if sec.get("subtitle"):
        textbox(s, 0.34, 1.62, 10.2, 0.43, f"■ {sec['subtitle']}",
                font=F_HEAD, size=18, bold=True, color=INK)
        content_top = 2.2
    return content_top


def build_body(prs, sec):
    """내지(본문) 슬라이드. 항상 '본문1'(빈 레이아웃) 베이스 + 코드 배치."""
    s = add_slide(prs, LAYOUTS["body"])
    content_top = _body_header(s, sec)
    content_bottom = 6.95
    area_h = content_bottom - content_top

    body = sec.get("body") or []
    table = sec.get("table")
    image = sec.get("image")
    columns = sec.get("columns")
    col_titles = sec.get("col_titles") or []

    LX, LW = 0.63, 4.61       # 좌측 컬럼
    RX, RW = 5.59, 4.61       # 우측 컬럼
    DIV_X = 5.41              # 컬럼 구분선

    def col_title(x, w, title):
        textbox(s, x, content_top, w, 0.35, title,
                font=F_TITLE, size=13, bold=False, color=HEADER)

    if columns and isinstance(columns, list) and len(columns) >= 2:
        # 2단 비교: 좌/우 글머리표 + 가운데 구분선
        ct = content_top
        if col_titles:
            col_title(LX, LW, col_titles[0] if len(col_titles) > 0 else "")
            if len(col_titles) > 1:
                col_title(RX, RW, col_titles[1])
            ct = content_top + 0.45
        vline(s, DIV_X, ct, content_bottom - ct, color=HEADER, weight=1.0)
        bullets(s, LX, ct, LW, content_bottom - ct, columns[0],
                font=F_BODY, size=11, color=GRAY)
        bullets(s, RX, ct, RW, content_bottom - ct, columns[1],
                font=F_BODY, size=11, color=GRAY)
    elif (table or image) and body:
        # 분할: 본문 좌측, 시각자료 우측
        ct = content_top
        if col_titles:
            col_title(LX, LW, col_titles[0] if len(col_titles) > 0 else "")
            if len(col_titles) > 1:
                col_title(RX, RW, col_titles[1])
            ct = content_top + 0.45
        bullets(s, LX, ct, LW, content_bottom - ct, body,
                font=F_BODY, size=11, color=GRAY)
        if table:
            add_table(s, RX, ct, RW, min(content_bottom - ct, 3.2),
                      table.get("headers"), table.get("rows", []))
        elif image:
            add_image_fit(s, _resolve(image, sec), RX, ct, RW,
                          content_bottom - ct)
    elif table:
        add_table(s, 0.45, content_top, 9.92, min(area_h, 3.6),
                  table.get("headers"), table.get("rows", []))
    elif image:
        add_image_fit(s, _resolve(image, sec), 0.45, content_top, 9.92, area_h)
    elif body:
        bullets(s, 0.63, content_top, 9.57, area_h, body,
                font=F_BODY, size=11, color=GRAY)

    # 출처
    if sec.get("source"):
        textbox(s, 0.45, 7.05, 9.9, 0.3, f"출처  I  {sec['source']}",
                font=F_SRC, size=9, color=GRAY)


def build_ending(prs, spec):
    s = add_slide(prs, LAYOUTS["ending"])
    text = spec.get("closing", "감사합니다.")
    textbox(s, 0.0, 3.29, SLIDE_W, 0.9, text,
            font=F_DIV, size=48, bold=False, color=NAVY,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)


# ---------------------------------------------------------------- main
def build(spec, out_path):
    prs = Presentation(TEMPLATE)
    wipe_slides(prs)

    build_cover(prs, spec)
    build_foreword(prs, spec)
    build_toc(prs, spec)

    for sec in spec.get("sections", []):
        if sec.get("divider"):
            build_divider(prs, sec["divider"])
        else:
            build_body(prs, sec)

    build_ending(prs, spec)
    prs.save(out_path)
    return out_path, len(prs.slides._sldIdLst)


def main():
    ap = argparse.ArgumentParser(description="사내 템플릿 기반 PPTX 보고서 생성")
    ap.add_argument("spec", help="deck 콘텐츠 JSON 파일 경로")
    ap.add_argument("-o", "--output", default="output.pptx", help="출력 PPTX 경로")
    args = ap.parse_args()

    if not os.path.exists(TEMPLATE):
        sys.exit(f"[오류] 템플릿을 찾을 수 없습니다: {TEMPLATE}")
    with open(args.spec, encoding="utf-8") as f:
        spec = json.load(f)

    global SPEC_DIR
    SPEC_DIR = os.path.dirname(os.path.abspath(args.spec))
    out, n = build(spec, args.output)
    print(f"[완료] {n}장 생성 -> {os.path.abspath(out)}")


if __name__ == "__main__":
    main()
